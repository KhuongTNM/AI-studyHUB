import logging
import pdfplumber
import docx
import pptx

logger = logging.getLogger(__name__)


def _clean_cell(cell) -> str:
    if cell is None:
        return ""
    return cell.strip().replace("\n", " ")


def extract_text_from_pdf(file_path: str) -> str:
    """
    Trích xuất text từ PDF, có xử lý riêng cho các trang chứa bảng
    (ví dụ bảng nhiều cột, danh sách dạng lưới) để tránh bị lộn xộn nội dung
    giữa các cột/hàng khi dùng extract_text() thông thường.
    """
    try:
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                tables = page.find_tables()

                if tables:
                    # 1. Loại bỏ vùng bảng ra khỏi phần text thường,
                    #    để không bị trích xuất trùng lặp hoặc lộn xộn.
                    non_table_page = page
                    for table in tables:
                        non_table_page = non_table_page.outside_bbox(table.bbox)

                    other_text = non_table_page.extract_text()
                    if other_text and other_text.strip():
                        text_parts.append(other_text.strip())

                    # 2. Trích xuất từng bảng theo đúng cấu trúc hàng/cột.
                    for table in tables:
                        data = table.extract()
                        if not data or len(data) < 2:
                            continue

                        header = [_clean_cell(c) for c in data[0]]

                        for row in data[1:]:
                            cells = [_clean_cell(c) for c in row]
                            parts = []
                            for h, c in zip(header, cells):
                                if not c:
                                    continue
                                parts.append(f"{h}: {c}" if h else c)
                            row_str = " | ".join(parts)
                            if row_str:
                                text_parts.append(row_str)
                else:
                    # Trang không có bảng -> trích xuất bình thường.
                    extracted = page.extract_text()
                    if extracted:
                        text_parts.append(extracted)

        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {e}")
        return ""


def _extract_docx_table(table) -> list[str]:
    """
    Trích xuất 1 bảng trong docx thành danh sách chuỗi,
    mỗi chuỗi là 1 hàng theo định dạng 'Header: Giá trị | Header: Giá trị'.
    Hàng đầu tiên được coi là header nếu không trống.
    """
    rows = []
    for cell in table.rows[0].cells:
        rows.append(cell.text.strip())
    header = rows if any(rows) else []

    result = []
    start = 1 if header else 0
    for row in table.rows[start:]:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        if not any(cells):
            continue
        if header:
            parts = []
            for h, c in zip(header, cells):
                if not c:
                    continue
                parts.append(f"{h}: {c}" if h else c)
            row_str = " | ".join(parts)
        else:
            row_str = " | ".join(c for c in cells if c)
        if row_str:
            result.append(row_str)
    return result


def extract_text_from_docx(file_path: str) -> str:
    """
    Trích xuất text từ DOCX, giữ đúng thứ tự xen kẽ giữa đoạn văn và bảng.
    Bảng được tách theo hàng/cột (giống xử lý PDF), không bị bỏ sót.
    """
    try:
        doc = docx.Document(file_path)
        text_parts = []

        # doc.element.body chứa tất cả phần tử theo đúng thứ tự trong file:
        # đoạn văn (w:p) xen kẽ bảng (w:tbl), đảm bảo không đọc nhầm thứ tự.
        from docx.oxml.ns import qn

        para_index = 0
        table_index = 0

        for child in doc.element.body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":  # đoạn văn
                if para_index < len(doc.paragraphs):
                    para = doc.paragraphs[para_index]
                    if para.text.strip():
                        text_parts.append(para.text.strip())
                    para_index += 1

            elif tag == "tbl":  # bảng
                if table_index < len(doc.tables):
                    table_rows = _extract_docx_table(doc.tables[table_index])
                    text_parts.extend(table_rows)
                    table_index += 1

        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = pptx.Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {e}")
        return ""


def extract_text(file_path: str, file_type: str) -> str:
    file_type = file_type.lower().strip('.')
    if file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_type in ['doc', 'docx']:
        return extract_text_from_docx(file_path)
    elif file_type in ['ppt', 'pptx']:
        return extract_text_from_pptx(file_path)
    else:
        logger.error(f"Unsupported file type: {file_type}")
        return ""
